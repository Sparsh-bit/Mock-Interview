"""
What a deck says — services/deck/extract.py

Text only. Rendering the slides to images is `render.py`, and the split matters: text
extraction is pure Python and always available, while rendering needs LibreOffice on the
host and is allowed to be absent.

THE BYTES HAVE ALREADY BEEN THROUGH `file_safety.verify` BY THE TIME ANYTHING HERE RUNS.
Nothing in this module decides what a file is; it is told, and it dispatches on the answer.
Reversing that — letting python-pptx "just try" the bytes — is the type-confusion hole that
`services/resume/file_safety.py` was written to close, and a deck upload is the same shape
of attack surface as a resume upload.
"""

from __future__ import annotations

import io

import structlog

from app.services.resume.file_safety import DocumentKind

logger = structlog.get_logger(__name__)

#: Deck text handed to the model, in characters.
#:
#: A 20-slide deck measures a few thousand; 40,000 is generous headroom for a text-heavy
#: appendix without letting a pathological file set the size of every downstream prompt.
#: Truncation is logged, because a silently shortened deck scores worse for a reason the
#: candidate cannot see.
MAX_DECK_CHARS = 40_000

#: Below this, there is nothing to judge.
#:
#: An image-only deck legitimately extracts almost no text, so this is NOT a refusal on its
#: own — the endpoint refuses only when rendering also produced nothing, because then there
#: is genuinely no evidence of any kind. See `DeckEvaluator`.
MIN_USEFUL_CHARS = 120


class DeckExtractionError(Exception):
    """Raised when a deck yields nothing that can be assessed."""

    def __init__(self, message: str, *, reason: str) -> None:
        super().__init__(message)
        self.reason = reason


def extract_text(data: bytes, kind: DocumentKind) -> str:
    """
    The deck's text, normalised, in reading order and capped.

    Returns "" rather than raising when a file of the right type simply has no text —
    that is a real and legitimate deck (all diagrams), and the caller decides whether the
    images make up for it.
    """
    if kind == "pptx":
        text = _pptx_text(data)
    elif kind == "pdf":
        text = _pdf_text(data)
    else:
        raise DeckExtractionError(
            f"A {kind} is not a deck.", reason="unsupported_kind"
        )

    text = _normalise(text)
    if len(text) > MAX_DECK_CHARS:
        logger.info("deck_text_truncated", chars=len(text), cap=MAX_DECK_CHARS)
        text = text[:MAX_DECK_CHARS]
    return text


def slide_count(data: bytes, kind: DocumentKind) -> int:
    """How many slides or pages, or 0 when it cannot be determined."""
    try:
        if kind == "pptx":
            from pptx import Presentation  # noqa: PLC0415 — optional at import time

            return len(Presentation(io.BytesIO(data)).slides)
        if kind == "pdf":
            import pypdf  # noqa: PLC0415

            return len(pypdf.PdfReader(io.BytesIO(data)).pages)
    except Exception as exc:  # noqa: BLE001 — a count is never worth failing an upload for
        logger.warning("deck_slide_count_failed", kind=kind, error=str(exc))
    return 0


def _pptx_text(data: bytes) -> str:
    """
    Every text frame, slide by slide, plus speaker notes.

    NOTES ARE INCLUDED AND LABELLED. A candidate who put their evidence in the notes has
    still done the work, and a judge reading the deck in the room would have them. Labelled
    so the model can tell a spoken aside from what is on the slide.

    TABLES TOO. python-pptx does not expose a table's cells through `shape.text_frame`, so
    a deck whose entire technical comparison is a table extracted as nothing — measured on
    a real deck, which is how this line came to exist.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []

    for index, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            slide_parts.extend(_shape_text(shape))
        if slide_parts:
            parts.append(f"--- Slide {index} ---")
            parts.extend(slide_parts)

        notes = _notes_text(slide)
        if notes:
            parts.append(f"[Speaker notes, slide {index}] {notes}")

    return "\n".join(parts)


def _shape_text(shape) -> list[str]:  # noqa: ANN001 — python-pptx shapes are untyped
    """Text out of one shape, including grouped shapes and table cells."""
    out: list[str] = []
    try:
        # A group's children are not reachable from `text_frame`, so recurse.
        if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
            for child in shape.shapes:
                out.extend(_shape_text(child))
            return out

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    out.append(" | ".join(cells))
            return out

        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                out.append(text)
    except Exception as exc:  # noqa: BLE001 — one unreadable shape is not a failed deck
        logger.debug("deck_shape_unreadable", error=str(exc))
    return out


def _notes_text(slide) -> str:  # noqa: ANN001
    try:
        if not slide.has_notes_slide:
            return ""
        return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def _pdf_text(data: bytes) -> str:
    """Page text via pypdf, which this backend already depends on for resumes."""
    import pypdf  # noqa: PLC0415

    reader = pypdf.PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception as exc:  # noqa: BLE001
            logger.debug("deck_pdf_page_unreadable", page=index, error=str(exc))
            continue
        if text:
            parts.append(f"--- Page {index} ---")
            parts.append(text)
    return "\n".join(parts)


def _normalise(text: str) -> str:
    """Collapse the whitespace a deck export produces, without losing line structure."""
    lines = [" ".join(line.split()) for line in (text or "").splitlines()]
    kept = [line for line in lines if line]
    return "\n".join(kept)
