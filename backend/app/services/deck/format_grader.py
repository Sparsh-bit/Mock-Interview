"""
Whether the deck is put together — services/deck/format_grader.py

Scores typography and layout out of 10, deterministically, with no model involved.

WHY THIS IS NOT ASKED OF THE MODEL. "Are the fonts consistent" is not a judgement, it is a
measurement: python-pptx can read every run's font name and size exactly. A language model
looking at rendered slides would estimate the same answer, less accurately, for the price
of a vision call — and it would give a different answer on the same deck twice. The other
eight criteria are judgements and belong to the model; this one does not.

PPTX AND PDF ARE GRADED ON DIFFERENT EVIDENCE, and the PDF path is honestly weaker. Font
metadata does not survive a PDF export, so there is no way to ask whether the title sizes
agree. The PDF checks are the structural ones that do survive: page geometry, and whether
anything is shaped like a heading. A deck uploaded as PDF is told which checks ran.
"""

from __future__ import annotations

import io
from collections import Counter
from dataclasses import dataclass, field

import structlog

from app.services.resume.file_safety import DocumentKind

logger = structlog.get_logger(__name__)

MAX_SCORE = 10

#: Typography floors, in points. A 24pt title and 16pt body is the widely-taught minimum
#: for a projected slide; below that a judge two rows back cannot read it.
_MIN_TITLE_PT = 24.0
_MIN_BODY_PT = 16.0

#: A deck should have a title on nearly every slide. 70% allows for a cover, a thank-you
#: slide and a full-bleed diagram without penalty.
_TITLE_COVERAGE = 0.70

#: Share of runs that must share one font family before it counts as consistent.
_DOMINANT_FONT_SHARE = 0.60

#: Distinct text colours past which a palette is not a palette.
_MAX_TEXT_COLOURS = 8

#: python-pptx's shape type for a picture.
_PICTURE = 13


@dataclass(slots=True)
class FormatReport:
    """A score out of `MAX_SCORE`, and the specific things behind it."""

    score: int = MAX_SCORE
    #: Candidate-facing observations. Deliberately actionable: each one names the fix.
    notes: list[str] = field(default_factory=list)
    #: Checks that could not run, so a low score is never mistaken for a failed parse.
    skipped: list[str] = field(default_factory=list)

    def penalise(self, points: float, note: str) -> None:
        self.score = max(0, self.score - int(round(points)))
        self.notes.append(note)


def grade(data: bytes, kind: DocumentKind, *, deck_text: str = "") -> FormatReport:
    """Grade a deck's formatting. Never raises — a parse failure is a reported skip."""
    try:
        if kind == "pptx":
            return _grade_pptx(data)
        if kind == "pdf":
            return _grade_pdf(data, deck_text)
    except Exception as exc:  # noqa: BLE001 — formatting is never worth failing an upload
        logger.warning("deck_format_grade_failed", kind=kind, error=str(exc))
        report = FormatReport()
        report.skipped.append("The formatting checks could not read this file.")
        return report

    report = FormatReport()
    report.skipped.append(f"Formatting is not graded for a {kind}.")
    return report


def _grade_pptx(data: bytes) -> FormatReport:
    from pptx import Presentation  # noqa: PLC0415
    from pptx.enum.shapes import PP_PLACEHOLDER  # noqa: PLC0415

    report = FormatReport()
    prs = Presentation(io.BytesIO(data))
    slides = list(prs.slides)

    if not slides:
        report.score = 0
        report.notes.append("The deck has no slides.")
        return report

    titled = 0
    fonts: list[str] = []
    title_pt: list[float] = []
    body_pt: list[float] = []
    colours: list[str] = []
    picture_slides = 0

    for slide in slides:
        if _has_title(slide, PP_PLACEHOLDER):
            titled += 1
        if any(getattr(s, "shape_type", None) == _PICTURE for s in slide.shapes):
            picture_slides += 1
        for shape in slide.shapes:
            _collect(shape, fonts, title_pt, body_pt, colours)

    # ── Titles ───────────────────────────────────────────────────────────────
    if titled / len(slides) < _TITLE_COVERAGE:
        report.penalise(
            2,
            f"Only {titled} of {len(slides)} slides use a title placeholder. "
            "A titled slide tells a reader what it is claiming before they read it.",
        )

    # ── One font family ──────────────────────────────────────────────────────
    if fonts:
        family, count = Counter(fonts).most_common(1)[0]
        if count / len(fonts) < _DOMINANT_FONT_SHARE:
            distinct = len(set(fonts))
            report.penalise(
                2,
                f"{distinct} font families are mixed across the deck, none dominant. "
                "Pick one for headings and one for body text.",
            )
        else:
            report.notes.append(f"Consistent typeface throughout ({family}).")
    else:
        report.skipped.append(
            "No font information is stored in this deck, so typeface consistency "
            "was not checked — the theme's defaults are being inherited."
        )

    # ── Legible sizes ────────────────────────────────────────────────────────
    small_titles = sum(1 for pt in title_pt if pt < _MIN_TITLE_PT)
    if small_titles:
        report.penalise(
            1,
            f"{small_titles} heading(s) are below {_MIN_TITLE_PT:.0f}pt.",
        )
    if body_pt:
        small_body = sum(1 for pt in body_pt if pt < _MIN_BODY_PT)
        if small_body / len(body_pt) > 0.30:
            report.penalise(
                2,
                f"{round(100 * small_body / len(body_pt))}% of body text is below "
                f"{_MIN_BODY_PT:.0f}pt, which does not project legibly.",
            )

    # ── Palette ──────────────────────────────────────────────────────────────
    distinct_colours = len(set(colours))
    if distinct_colours > _MAX_TEXT_COLOURS:
        report.penalise(
            1,
            f"{distinct_colours} distinct text colours are in use. "
            "Tighten it to a palette of three or four.",
        )

    # ── Anything to look at ──────────────────────────────────────────────────
    if picture_slides == 0:
        report.penalise(
            2,
            "No images or diagrams anywhere in the deck. A technical idea is "
            "almost always faster to show than to describe.",
        )

    return report


def _has_title(slide, placeholder_enum) -> bool:  # noqa: ANN001
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type in (
                placeholder_enum.TITLE,
                placeholder_enum.CENTER_TITLE,
            ):
                # An EMPTY title placeholder is not a title. The layout puts one on
                # every slide whether or not anybody typed in it, so counting the
                # placeholder rather than its contents scored blank decks full marks.
                return bool((shape.text_frame.text or "").strip())
        except Exception:  # noqa: BLE001
            continue
    return False


def _collect(  # noqa: ANN001
    shape,
    fonts: list[str],
    title_pt: list[float],
    body_pt: list[float],
    colours: list[str],
) -> None:
    """Font facts out of one shape, recursing into groups."""
    try:
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                _collect(child, fonts, title_pt, body_pt, colours)
            return
        if not getattr(shape, "has_text_frame", False):
            return

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = getattr(run, "font", None)
                if font is None:
                    continue

                name = (getattr(font, "name", None) or "").strip()
                if name:
                    fonts.append(name)

                size = getattr(getattr(font, "size", None), "pt", None)
                if size:
                    size = float(size)
                    is_heading = paragraph.level == 0 and (
                        bool(getattr(font, "bold", False)) or size >= _MIN_TITLE_PT
                    )
                    (title_pt if is_heading else body_pt).append(size)

                rgb = getattr(getattr(font, "color", None), "rgb", None)
                if rgb is not None:
                    colours.append(str(rgb))
    except Exception as exc:  # noqa: BLE001
        logger.debug("deck_format_shape_unreadable", error=str(exc))


def _grade_pdf(data: bytes, deck_text: str) -> FormatReport:
    import pypdf  # noqa: PLC0415

    report = FormatReport()
    report.skipped.append(
        "This was uploaded as a PDF, so font sizes and typeface consistency could not "
        "be checked — that information does not survive a PDF export. Upload the .pptx "
        "for the full formatting review."
    )

    reader = pypdf.PdfReader(io.BytesIO(data))
    pages = reader.pages
    if not pages:
        report.score = 0
        report.notes.append("The file has no pages.")
        return report

    # ── Consistent page geometry ─────────────────────────────────────────────
    sizes = {
        (round(float(p.mediabox.width)), round(float(p.mediabox.height))) for p in pages
    }
    if len(sizes) > 1:
        report.penalise(
            2, f"{len(sizes)} different page sizes in one file. Export at a single size."
        )

    # ── Landscape, like a projector ──────────────────────────────────────────
    first = pages[0].mediabox
    if float(first.height) > float(first.width):
        report.penalise(
            1,
            "The pages are portrait. A deck shown on a projector or a shared screen "
            "should be landscape.",
        )

    # ── Something shaped like headings ───────────────────────────────────────
    headings = [
        line
        for line in (deck_text or "").splitlines()
        if len(line.strip()) >= 6 and line.strip() == line.strip().upper()
    ]
    if len(headings) < 3:
        report.penalise(
            2,
            "Few lines read as headings. Give every slide a short title line so a "
            "reader can skim the argument.",
        )

    return report
