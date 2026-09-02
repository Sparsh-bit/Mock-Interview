"""
The deck review — tests/test_deck_review.py

Covers the four things that would each be a silent wrong answer rather than an error:

  1. A .pptx is recognised BY ITS BYTES, goes through the same archive and macro checks a
     .docx does, and is NOT accepted as a resume. The last one is the regression that
     adding `pptx` to `DocumentKind` invites: one shared `verify` serving two endpoints
     that accept different things.

  2. A request carrying images never reaches a provider that cannot see them. There is no
     error when it does — the model answers normally, having scored the deck on text it was
     given and images it silently dropped — so nothing but a test notices.

  3. The rubric's weights total 100. The rubric this was ported from summed to 105, so a
     perfect deck scored 105.00 and every percentage the product displayed was of nothing.

  4. Formatting is MEASURED, not asked of the model. The model's answer for that criterion
     is overwritten, so a model that returns nine scores cannot beat the parser.

WHAT IS NOT HERE. Nothing calls a live model. The judging pass is faked at the provider
boundary, which is the same seam `generate_structured` retries against, so the wiring under
test is the real wiring.
"""

from __future__ import annotations

import io
import uuid
import zipfile

import pytest
from httpx import ASGITransport, AsyncClient
from jose import jwt
from pydantic import BaseModel

from app.core.config import settings
from app.main import app
from app.services.ai.base_provider import (
    BaseAIProvider,
    ImagePart,
    ProviderMessage,
    ProviderRequest,
    ProviderResponse,
)
from app.services.billing.plans import FEATURES, TRIAL_ALLOWANCE, items_for
from app.services.deck import criteria, extract, format_grader
from app.services.resume.file_safety import (
    DECK_KINDS,
    RESUME_KINDS,
    UnsafeDocument,
    sniff,
    verify,
)


class _Scores(BaseModel):
    """The shape the faked provider answers with, for the vision-gate tests."""

    scores: dict[str, int] = {}
    summary: str = ""


# ─── Fixtures built in-test, so the suite depends on no binary assets ────────


def _tiny_png() -> bytes:
    """A valid 8x8 PNG, hand-assembled to avoid needing an image library."""
    import struct
    import zlib

    w = h = 8
    raw = b"".join(b"\x00" + bytes(range(w * 3)) for _ in range(h))

    def chunk(tag: bytes, data: bytes) -> bytes:
        body = tag + data
        return (
            struct.pack(">I", len(data))
            + body
            + struct.pack(">I", zlib.crc32(body) & 0xFFFFFFFF)
        )

    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )


def build_pptx(*, slides: int = 3, tidy: bool = True) -> bytes:
    """A real .pptx. `tidy=False` breaks every formatting rule the grader checks."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    fonts = ["Arial", "Times New Roman", "Comic Sans MS", "Courier New", "Verdana", "Georgia"]

    for i in range(slides):
        if tidy:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Section {i + 1}"
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.name = "Calibri"
                run.font.size = Pt(32)
                run.font.bold = True
            frame = slide.placeholders[1].text_frame
            frame.text = (
                "Problem statement and target user persona. Architecture: FastAPI, "
                "Docker, Postgres. Latency 40ms p95, accuracy 94% F1. TAM $2B."
            )
            for para in frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(18)
            slide.shapes.add_picture(
                io.BytesIO(_tiny_png()), Inches(6.0), Inches(5.0), Inches(2.0), Inches(1.2)
            )
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank: no title
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(4))
            frame = box.text_frame
            frame.text = f"words {i}"
            for j in range(4):
                para = frame.add_paragraph()
                para.text = f"small line {j}"
                for run in para.runs:
                    run.font.name = fonts[(i + j) % len(fonts)]
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(
                        (i * 40) % 256, (j * 70) % 256, (i * j * 30) % 256
                    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pdf() -> bytes:
    """A one-page PDF with extractable, heading-shaped text."""
    from tests.pdf_builder import build_pdf as _build
    from tests.pdf_builder import visible_run

    return _build(
        visible_run(
            "PROBLEM STATEMENT Our users lose time. ARCHITECTURE FastAPI and Postgres. "
            "MARKET TAM is large. IMPACT measured in hours saved."
        )
    )


def with_member(data: bytes, name: str, payload: bytes) -> bytes:
    """The same archive with one extra member — how a macro gets into a deck."""
    buf = io.BytesIO(data)
    with zipfile.ZipFile(buf, "a") as archive:
        archive.writestr(name, payload)
    return buf.getvalue()


# ─── 1. What the bytes are ───────────────────────────────────────────────────


class TestWhatTheBytesAre:
    def test_a_pptx_is_recognised_from_its_bytes(self):
        assert sniff(build_pptx()) == "pptx"

    def test_a_plain_zip_is_not_a_deck(self):
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as archive:
            archive.writestr("hello.txt", b"not a presentation")
        assert sniff(buf.getvalue()) is None

    def test_an_ooxml_file_that_is_not_a_presentation_is_not_a_deck(self):
        """A .docx shares the signature and the manifest. Only the main part differs."""
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as archive:
            archive.writestr("[Content_Types].xml", b"<Types/>")
            archive.writestr("word/document.xml", b"<document/>")
        assert sniff(buf.getvalue()) == "docx"

    def test_the_deck_endpoint_accepts_a_presentation(self):
        assert verify(build_pptx(), allowed=DECK_KINDS) == "pptx"

    def test_the_resume_endpoint_still_refuses_a_presentation(self):
        """
        THE REGRESSION THAT ADDING pptx INVITES. One `verify` serves both uploads, so
        teaching it a new format could silently make a slide deck a valid CV.
        """
        with pytest.raises(UnsafeDocument) as caught:
            verify(build_pptx(), allowed=RESUME_KINDS)
        assert caught.value.reason == "unsupported_type"
        assert "PowerPoint" in str(caught.value)

    def test_the_deck_endpoint_refuses_a_word_document(self):
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as archive:
            archive.writestr("[Content_Types].xml", b"<Types/>")
            archive.writestr("word/document.xml", b"<document/>")
        with pytest.raises(UnsafeDocument):
            verify(buf.getvalue(), allowed=DECK_KINDS)

    def test_a_macro_bearing_deck_is_refused(self):
        """A .pptx is the same zip container, so it gets the same member scan."""
        armed = with_member(build_pptx(), "ppt/vbaProject.bin", b"\xd0\xcf\x11\xe0payload")
        with pytest.raises(UnsafeDocument) as caught:
            verify(armed, allowed=DECK_KINDS)
        assert caught.value.reason == "pptx_macro"

    def test_a_deck_with_an_embedded_object_is_refused(self):
        armed = with_member(build_pptx(), "ppt/embeddings/oleObject1.bin", b"payload")
        with pytest.raises(UnsafeDocument) as caught:
            verify(armed, allowed=DECK_KINDS)
        assert caught.value.reason == "pptx_macro"

    def test_the_refusal_reason_names_the_format(self):
        """
        `docx_macro` is asserted verbatim by test_upload_file_safety, so the reason had to
        stay per-kind rather than become one shared string.
        """
        armed = with_member(build_pptx(), "ppt/vbaProject.bin", b"x")
        with pytest.raises(UnsafeDocument) as caught:
            verify(armed, allowed=DECK_KINDS)
        assert caught.value.reason.startswith("pptx_")


# ─── 2. Vision cannot be dropped silently ────────────────────────────────────


class _Blind(BaseAIProvider):
    """A text-only provider. Answers perfectly well, having seen nothing."""

    saw_images = False

    @property
    def provider_name(self) -> str:
        return "blind"

    @property
    def model_name(self) -> str:
        return "text-only-1"

    async def complete(self, request: ProviderRequest) -> ProviderResponse:
        # RECORDED ON `_Blind` DELIBERATELY, not on `type(self)`. `_Seeing` subclasses
        # this, so `type(self)` would set the flag on the subclass and leave the one the
        # tests read untouched — a green test that checked nothing.
        _Blind.saw_images = request.has_images
        return ProviderResponse(
            content='{"scores": {"innovation": 7}, "summary": "fine"}',
            model=self.model_name,
            prompt_tokens=1,
            completion_tokens=1,
            finish_reason="stop",
        )

    async def health_check(self) -> bool:
        return True


class _Seeing(_Blind):
    @property
    def provider_name(self) -> str:
        return "seeing"

    @property
    def supports_vision(self) -> bool:
        return True


class TestVisionIsNeverDroppedSilently:
    """
    A blind provider handed images does not fail — it answers. That is the whole problem:
    the score comes back looking exactly like a real one.
    """

    @pytest.fixture(autouse=True)
    def _reset(self):
        _Blind.saw_images = False
        yield

    @pytest.mark.asyncio
    async def test_an_image_request_skips_a_blind_provider(self, monkeypatch):
        from app.core.exceptions import AIProviderUnavailableError
        from app.services.ai import generate as generate_mod

        monkeypatch.setattr(generate_mod, "get_ai_providers", lambda: [_Blind()])

        with pytest.raises(AIProviderUnavailableError):
            await generate_mod.generate_structured(
                _Scores,
                [
                    ProviderMessage(
                        role="user",
                        content="score it",
                        images=(ImagePart(base64_data="QUJD"),),
                    )
                ],
                max_tokens=64,
                context="deck_evaluation",
            )
        assert _Blind.saw_images is False, "the blind provider was called with images"

    @pytest.mark.asyncio
    async def test_a_text_only_request_still_uses_a_blind_provider(self, monkeypatch):
        """The gate must not touch the other calls in the application."""
        from app.services.ai import generate as generate_mod

        monkeypatch.setattr(generate_mod, "get_ai_providers", lambda: [_Blind()])
        parsed, _ = await generate_mod.generate_structured(
            _Scores,
            [ProviderMessage(role="user", content="score it")],
            max_tokens=64,
            context="deck_evaluation",
        )
        assert parsed.scores == {"innovation": 7}

    @pytest.mark.asyncio
    async def test_a_seeing_provider_receives_the_images(self, monkeypatch):
        from app.services.ai import generate as generate_mod

        monkeypatch.setattr(generate_mod, "get_ai_providers", lambda: [_Seeing()])
        await generate_mod.generate_structured(
            _Scores,
            [
                ProviderMessage(
                    role="user", content="score it", images=(ImagePart(base64_data="QUJD"),)
                )
            ],
            max_tokens=64,
            context="deck_evaluation",
        )
        assert _Blind.saw_images is True


class TestTheProvidersTranslateImages:
    def test_openai_shape_keeps_text_only_messages_as_strings(self):
        """
        Every existing call is text-only and has always sent a bare string. Changing that
        shape would change the cache key of calls that have nothing to do with vision.
        """
        from app.services.ai.glm_provider import _to_api_messages

        out = _to_api_messages([ProviderMessage(role="user", content="hi")])
        assert out == [{"role": "user", "content": "hi"}]

    def test_openai_shape_emits_image_url_parts(self):
        from app.services.ai.glm_provider import _to_api_messages

        out = _to_api_messages(
            [
                ProviderMessage(
                    role="user",
                    content="look",
                    images=(ImagePart(base64_data="QUJD", media_type="image/png"),),
                )
            ]
        )
        parts = out[0]["content"]
        assert parts[0] == {"type": "text", "text": "look"}
        assert parts[1]["image_url"]["url"] == "data:image/png;base64,QUJD"

    def test_no_unknown_keys_reach_the_api(self):
        """
        The payload was `m.model_dump()`, which started emitting `images` the moment the
        field existed — an unknown key these servers answer with a 400.
        """
        from app.services.ai.glm_provider import _to_api_messages

        for message in _to_api_messages([ProviderMessage(role="user", content="hi")]):
            assert set(message) == {"role", "content"}

    def test_anthropic_puts_the_text_after_the_images(self):
        from app.services.ai.anthropic_provider import _turn_content

        blocks = _turn_content(
            ProviderMessage(
                role="user", content="score it", images=(ImagePart(base64_data="QUJD"),)
            )
        )
        assert [b["type"] for b in blocks] == ["image", "text"]
        assert blocks[0]["source"] == {
            "type": "base64",
            "media_type": "image/jpeg",
            "data": "QUJD",
        }

    def test_anthropic_keeps_text_only_turns_as_strings(self):
        from app.services.ai.anthropic_provider import _turn_content

        assert _turn_content(ProviderMessage(role="user", content="hi")) == "hi"

    def test_an_unsupported_media_type_is_refused_at_the_edge(self):
        with pytest.raises(ValueError):
            ImagePart(base64_data="x", media_type="image/tiff")  # type: ignore[arg-type]


class TestImagesAreTreatedAsUntrusted:
    def test_the_fence_rule_is_attached_when_images_are_present(self):
        """
        Text on a slide arrives as pixels, where no delimiter can wrap it. The system block
        must still carry the data-versus-instruction framing.
        """
        from app.prompts.prompt_loader import get_prompt_loader
        from app.services.ai.prompt_builder import PromptBuilder
        from app.services.ai.untrusted import FENCE_RULE

        messages = PromptBuilder(get_prompt_loader()).chat(
            system_template="deck_diagrams",
            user_content="Analyse the attached slides.",
            images=(ImagePart(base64_data="QUJD"),),
        )
        assert FENCE_RULE in messages[0].content

    def test_the_vision_prompt_says_text_in_a_picture_is_not_an_instruction(self):
        from app.prompts.prompt_loader import get_prompt_loader

        for name in ("deck_diagrams", "deck_evaluator"):
            body = get_prompt_loader().load(name).lower()
            assert "instruction" in body
            assert "image" in body


# ─── 3. The rubric ───────────────────────────────────────────────────────────


class TestTheRubric:
    def test_the_weights_total_one_hundred(self):
        """
        THE PORTED RUBRIC SUMMED TO 105. A perfect deck scored 105.00 and every percentage
        shown was of nothing. This is the pin that stops it coming back.
        """
        assert criteria.total_weight() == 100

    def test_a_perfect_deck_scores_exactly_one_hundred(self):
        perfect = {c.key: c.max_score for c in criteria.CRITERIA}
        assert criteria.weighted_total(perfect) == 100.0

    def test_an_empty_score_map_is_zero_not_an_error(self):
        assert criteria.weighted_total({}) == 0.0

    def test_out_of_range_scores_are_clamped(self):
        clamped = criteria.clamp_scores({"innovation": 99, "impact": -4})
        assert clamped == {"innovation": 10, "impact": 1}

    def test_unknown_criteria_are_dropped(self):
        assert criteria.clamp_scores({"vibes": 10}) == {}

    def test_non_numeric_scores_are_dropped_not_crashed(self):
        assert criteria.clamp_scores({"innovation": "excellent"}) == {}

    def test_every_criterion_appears_in_the_prompt_block(self):
        block = criteria.criteria_block()
        for criterion in criteria.CRITERIA:
            assert criterion.key in block


# ─── 4. Formatting is measured, not judged ───────────────────────────────────


class TestTheFormatGrader:
    def test_a_tidy_deck_scores_well(self):
        report = format_grader.grade(build_pptx(tidy=True), "pptx")
        assert report.score >= 8, report.notes

    def test_a_sloppy_deck_scores_badly(self):
        report = format_grader.grade(build_pptx(tidy=False), "pptx")
        assert report.score <= 3, report.notes

    def test_it_names_what_is_wrong(self):
        notes = " ".join(format_grader.grade(build_pptx(tidy=False), "pptx").notes).lower()
        assert "title" in notes
        assert "font" in notes

    def test_an_empty_title_placeholder_is_not_a_title(self):
        """
        The layout puts a title placeholder on every slide whether or not anybody typed in
        it, so counting placeholders rather than their contents scored blank decks full
        marks.
        """
        from pptx import Presentation

        prs = Presentation()
        for _ in range(4):
            prs.slides.add_slide(prs.slide_layouts[1])  # title layout, nothing typed
        buf = io.BytesIO()
        prs.save(buf)
        notes = " ".join(format_grader.grade(buf.getvalue(), "pptx").notes).lower()
        assert "title" in notes

    def test_a_pdf_says_which_checks_could_not_run(self):
        report = format_grader.grade(build_pdf(), "pdf", deck_text="PROBLEM\nARCHITECTURE")
        assert any("pdf" in s.lower() for s in report.skipped)

    def test_unreadable_bytes_are_a_reported_skip_not_an_exception(self):
        report = format_grader.grade(b"not a deck at all", "pptx")
        assert report.skipped
        assert 0 <= report.score <= format_grader.MAX_SCORE

    def test_the_score_never_leaves_its_range(self):
        for data, kind in ((build_pptx(tidy=False), "pptx"), (build_pdf(), "pdf")):
            report = format_grader.grade(data, kind)
            assert 0 <= report.score <= format_grader.MAX_SCORE


class TestTheMeasuredScoreBeatsTheModel:
    @pytest.mark.asyncio
    async def test_the_format_criterion_comes_from_the_parser(self, monkeypatch):
        """
        The model is not asked for this criterion, but a model asked for nine scores will
        occasionally return nine. Whatever it says about formatting is a guess about
        something a parser already measured exactly.
        """
        from app.services.deck import evaluator as evaluator_mod

        deck = build_pptx(tidy=False)
        measured = format_grader.grade(deck, "pptx").score

        async def _fake_judge(self, deck_text, diagrams, rendered, log):  # noqa: ANN001
            from app.services.deck.schemas import JudgeResponse

            return JudgeResponse(
                scores={c.key: 10 for c in criteria.CRITERIA}, summary="all tens"
            )

        monkeypatch.setattr(evaluator_mod.DeckEvaluator, "_judge", _fake_judge)
        monkeypatch.setattr(
            evaluator_mod.DeckEvaluator, "_vision_available", staticmethod(lambda: False)
        )

        result = await evaluator_mod.DeckEvaluator().evaluate(deck, "pptx", filename="d.pptx")
        row = next(r for r in result.scores if r.key == criteria.FORMAT_CRITERION)
        assert row.score == measured
        assert row.score != 10, "the model's guess overwrote the measurement"
        assert row.measured is True


# ─── Extraction ──────────────────────────────────────────────────────────────


class TestExtraction:
    def test_slide_text_is_extracted_in_order(self):
        text = extract.extract_text(build_pptx(slides=3), "pptx")
        assert "Section 1" in text
        assert text.index("Section 1") < text.index("Section 3")

    def test_the_slide_count_is_read(self):
        assert extract.slide_count(build_pptx(slides=5), "pptx") == 5

    def test_speaker_notes_are_included_and_labelled(self):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"
        slide.notes_slide.notes_text_frame.text = "the evidence is in the notes"
        buf = io.BytesIO()
        prs.save(buf)

        text = extract.extract_text(buf.getvalue(), "pptx")
        assert "the evidence is in the notes" in text
        assert "Speaker notes" in text

    def test_table_cells_are_extracted(self):
        """
        python-pptx does not expose table cells through `text_frame`, so a deck whose whole
        comparison is a table extracted as nothing.
        """
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
        table.cell(0, 0).text = "Latency"
        table.cell(0, 1).text = "40ms"
        buf = io.BytesIO()
        prs.save(buf)

        text = extract.extract_text(buf.getvalue(), "pptx")
        assert "Latency" in text
        assert "40ms" in text

    def test_extraction_is_capped(self, monkeypatch):
        monkeypatch.setattr(extract, "MAX_DECK_CHARS", 50)
        assert len(extract.extract_text(build_pptx(slides=6), "pptx")) <= 50

    def test_a_non_deck_kind_is_refused(self):
        with pytest.raises(extract.DeckExtractionError):
            extract.extract_text(b"x", "docx")


# ─── Billing ─────────────────────────────────────────────────────────────────


class TestItIsMetered:
    def test_deck_is_a_metered_feature(self):
        assert "deck" in FEATURES

    def test_it_has_an_explicit_trial_decision(self):
        assert TRIAL_ALLOWANCE["deck"] == 0

    def test_it_can_actually_be_bought(self):
        items = items_for("deck")
        assert items, "a metered feature nobody can buy is a paywall with no door"
        assert all(i.price_paise > 0 for i in items)

    def test_both_ai_call_sites_bill_to_it(self):
        """
        `scripts/item_margin.py` maps a `context=` label to a billable feature. A label
        missing from it is not a runtime error — that feature's AI cost simply appears in no
        margin line, and the item reads as pure profit.
        """
        from scripts.item_margin import _AI_FEATURE_TO_BILLABLE

        assert _AI_FEATURE_TO_BILLABLE["deck_diagrams"] == "deck"
        assert _AI_FEATURE_TO_BILLABLE["deck_evaluation"] == "deck"


# ─── The endpoint ────────────────────────────────────────────────────────────


def _token(user_id: uuid.UUID) -> str:
    return jwt.encode(
        {
            "sub": str(user_id),
            "email": f"deck-{user_id}@example.test",
            "aud": settings.SUPABASE_JWT_AUDIENCE,
            "role": "authenticated",
        },
        settings.SUPABASE_JWT_SECRET,
        algorithm="HS256",
    )


class TestTheEndpoint:
    """
    The order of refusals, which is the part worth pinning: consent, then size, then what
    the bytes are, and only then the charge.
    """

    @pytest.fixture
    async def user_without_consent(self):
        from app.db.session import AsyncSessionFactory, engine
        from app.models.base import Base
        from app.models.user import User

        async with engine.begin() as conn:
            await conn.run_sync(Base.metadata.create_all)
        uid = uuid.uuid4()
        async with AsyncSessionFactory() as db:
            db.add(
                User(
                    id=uid,
                    supabase_uid=str(uid),
                    email=f"deck-{uid}@example.test",
                    is_active=True,
                    is_admin=False,
                )
            )
            await db.commit()
        return uid

    @pytest.fixture
    async def user(self, user_without_consent):
        from app.db.session import AsyncSessionFactory
        from app.models.consent import PURPOSE_DECK_PROCESSING, SOURCE_SIGNUP
        from app.services.legal.consent import record as record_consent

        async with AsyncSessionFactory() as db:
            await record_consent(
                db,
                user_without_consent,
                purpose=PURPOSE_DECK_PROCESSING,
                granted=True,
                source=SOURCE_SIGNUP,
            )
            await db.commit()
        return user_without_consent

    async def _post(self, uid: uuid.UUID, name: str, data: bytes, content_type: str):
        transport = ASGITransport(app=app, raise_app_exceptions=False)
        async with (
            app.router.lifespan_context(app),
            AsyncClient(transport=transport, base_url="http://test", timeout=60.0) as ac,
        ):
            return await ac.post(
                "/api/v1/deck/review",
                headers={"Authorization": f"Bearer {_token(uid)}"},
                files={"file": (name, data, content_type)},
            )

    _PPTX_MIME = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    @pytest.mark.asyncio
    async def test_it_requires_authentication(self):
        transport = ASGITransport(app=app, raise_app_exceptions=False)
        async with (
            app.router.lifespan_context(app),
            AsyncClient(transport=transport, base_url="http://test", timeout=30.0) as ac,
        ):
            response = await ac.post(
                "/api/v1/deck/review",
                files={"file": ("d.pptx", build_pptx(), self._PPTX_MIME)},
            )
        assert response.status_code in {401, 403}

    @pytest.mark.asyncio
    async def test_consent_is_required_before_the_bytes_are_touched(
        self, user_without_consent
    ):
        response = await self._post(
            user_without_consent, "d.pptx", build_pptx(), self._PPTX_MIME
        )
        assert response.status_code == 428
        assert response.json()["detail"]["purpose"] == "deck_processing"

    @pytest.mark.asyncio
    async def test_an_empty_file_is_refused(self, user):
        assert (await self._post(user, "d.pptx", b"", self._PPTX_MIME)).status_code == 422

    @pytest.mark.asyncio
    async def test_an_oversized_deck_is_refused(self, user):
        oversized = b"\x00" * (settings.deck_upload_size_bytes + 1)
        assert (await self._post(user, "d.pptx", oversized, self._PPTX_MIME)).status_code == 413

    @pytest.mark.asyncio
    async def test_bytes_that_are_not_a_deck_are_refused_even_when_labelled_one(self, user):
        response = await self._post(user, "d.pptx", b"just some text", self._PPTX_MIME)
        assert response.status_code == 415

    @pytest.mark.asyncio
    async def test_a_macro_bearing_deck_is_refused_by_the_endpoint(self, user):
        armed = with_member(build_pptx(), "ppt/vbaProject.bin", b"payload")
        assert (await self._post(user, "d.pptx", armed, self._PPTX_MIME)).status_code == 415

    @pytest.fixture
    async def paid_user(self, user):
        """A user with one deck review to spend, so the happy path can actually run."""
        from app.db.session import AsyncSessionFactory
        from app.services.billing.credits import grant

        async with AsyncSessionFactory() as db:
            # PER-USER `payment_ref`, BECAUSE `grant` IS IDEMPOTENT ON IT. A shared literal
            # made the first test's grant land and every later one a silent no-op — the
            # tables are created once for the run, not truncated between tests, so the
            # second call saw a ref it had already applied and returned False. The symptom
            # was a 402 in a test whose whole point was the happy path.
            await grant(
                db, user, "deck", 1, kind="purchase", payment_ref=f"test-deck-{user}"
            )
            await db.commit()
        return user

    @pytest.mark.asyncio
    async def test_a_paid_review_returns_a_complete_evaluation(self, paid_user, monkeypatch):
        """
        THE HAPPY PATH, OVER HTTP, WITH THE MODEL FAKED AT THE PROVIDER BOUNDARY.

        Everything between the multipart body and the provider is the real thing: file
        safety, the credit charge, extraction, the format parser, prompt assembly and the
        response schema. Only the two model calls are stood in for, which is the one part
        that cannot run in a unit suite.

        Vision is switched off here rather than faked. A seeing provider would need the
        renderer, which needs LibreOffice, which is deliberately absent from CI — so this
        asserts the DEGRADED path reports itself honestly, which is the case that will
        actually run on a default deployment.
        """
        from app.services.ai import generate as generate_mod
        from app.services.deck import evaluator as evaluator_mod

        monkeypatch.setattr(generate_mod, "get_ai_providers", lambda: [_Blind()])
        monkeypatch.setattr(
            evaluator_mod.DeckEvaluator, "_vision_available", staticmethod(lambda: False)
        )

        response = await self._post(paid_user, "pitch.pptx", build_pptx(), self._PPTX_MIME)
        assert response.status_code == 200, response.text
        body = response.json()

        # Every criterion is present and in range, in the rubric's own order.
        assert [row["key"] for row in body["scores"]] == [
            c.key for c in criteria.CRITERIA
        ]
        for row in body["scores"]:
            assert 0 <= row["score"] <= row["max_score"], row

        # The total is a percentage of a rubric that sums to 100, so it cannot exceed it.
        assert 0.0 <= body["weighted_total"] <= 100.0

        assert body["filename"] == "pitch.pptx"
        assert body["slide_count"] == 3

        # DEGRADED, AND IT SAYS SO. A text-only score presented as a complete one is the
        # failure this field exists to prevent.
        assert body["images_analysed"] == 0
        assert body["vision_unavailable_reason"] == "no_vision_provider"

    @pytest.mark.asyncio
    async def test_the_review_is_charged_exactly_once(self, paid_user, monkeypatch):
        """One granted review buys one review, and the second attempt is refused."""
        from app.services.ai import generate as generate_mod
        from app.services.deck import evaluator as evaluator_mod

        monkeypatch.setattr(generate_mod, "get_ai_providers", lambda: [_Blind()])
        monkeypatch.setattr(
            evaluator_mod.DeckEvaluator, "_vision_available", staticmethod(lambda: False)
        )

        first = await self._post(paid_user, "a.pptx", build_pptx(), self._PPTX_MIME)
        assert first.status_code == 200, first.text

        second = await self._post(paid_user, "b.pptx", build_pptx(), self._PPTX_MIME)
        assert second.status_code == 402

    @pytest.mark.asyncio
    async def test_an_unreadable_deck_does_not_keep_the_charge(self, paid_user, monkeypatch):
        """
        A 422 must leave the balance untouched. `consume` does not commit — `get_db` rolls
        the request back on the exception — so the refusal and the refund are the same act.
        Asserted by spending the credit afterwards: if the failed attempt had banked it,
        this second call would be a 402.
        """
        from app.services.ai import generate as generate_mod
        from app.services.deck import evaluator as evaluator_mod

        monkeypatch.setattr(generate_mod, "get_ai_providers", lambda: [_Blind()])
        monkeypatch.setattr(
            evaluator_mod.DeckEvaluator, "_vision_available", staticmethod(lambda: False)
        )

        # A structurally valid .pptx with no text in it at all, and no images to fall back
        # on, is the "nothing to assess" case.
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        empty = io.BytesIO()
        prs.save(empty)

        refused = await self._post(paid_user, "empty.pptx", empty.getvalue(), self._PPTX_MIME)
        assert refused.status_code == 422, refused.text

        # The credit survived the refusal.
        after = await self._post(paid_user, "real.pptx", build_pptx(), self._PPTX_MIME)
        assert after.status_code == 200, after.text

    @pytest.mark.asyncio
    async def test_an_exhausted_balance_is_refused_before_any_ai_call(self, user):
        """
        The trial allowance is zero, so a fresh account has nothing to spend. 402 here also
        proves the charge happens before the evaluation rather than after it.
        """
        response = await self._post(user, "d.pptx", build_pptx(), self._PPTX_MIME)
        assert response.status_code == 402
